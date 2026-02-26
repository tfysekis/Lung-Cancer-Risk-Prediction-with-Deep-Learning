"""
Create a PowerPoint presentation from project results and visuals.
Run from project root:  python src/create_presentation.py
Output: reports/Lung_Cancer_Deep_Learning_Presentation.pptx
"""

import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError as e:
    print("Failed to import python-pptx:", e)
    print("Install with:  pip install python-pptx")
    sys.exit(1)

# Paths
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.dirname(SCRIPT_DIR)
REPORTS = os.path.join(PROJECT_ROOT, "reports")
SINGLE = os.path.join(REPORTS, "single_split")
CV5 = os.path.join(REPORTS, "cv5")
CV10 = os.path.join(REPORTS, "cv10")
OUTPUT_PATH = os.path.join(REPORTS, "Lung_Cancer_Deep_Learning_Presentation.pptx")

# All 7 architectures: (file_suffix, display_name) — must match compare_models output
MODELS = [
    ("standard_128-64-32", "Standard (128-64-32)"),
    ("simple_64-32", "Simple (64-32)"),
    ("minimal_32", "Minimal (32)"),
    ("deep_5_layers", "Deep (5 layers)"),
    ("elu_activation", "ELU Activation"),
    ("layernorm", "LayerNorm"),
    ("residual", "Residual"),
]


def add_title_slide(prs, title, subtitle=""):
    layout = prs.slide_layouts[0]  # title slide
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
    return slide


def add_content_slide(prs, title, bullets):
    layout = prs.slide_layouts[1]  # title and content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for b in bullets:
        p = body.add_paragraph()
        p.text = b
        p.level = 0
    return slide


def add_slide_with_image(prs, title, image_path, caption=""):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    # Title
    left = Inches(0.5)
    top = Inches(0.3)
    w, h = Inches(9), Inches(0.6)
    tx = slide.shapes.add_textbox(left, top, w, h)
    tx.text_frame.text = title
    tx.text_frame.paragraphs[0].font.size = Pt(24)
    tx.text_frame.paragraphs[0].font.bold = True
    # Image (scale to fit on slide)
    if os.path.isfile(image_path):
        slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.0), width=Inches(9), height=Inches(5.2))
    if caption:
        tx2 = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.4))
        tx2.text_frame.text = caption
        tx2.text_frame.paragraphs[0].font.size = Pt(12)
    return slide


def add_slide_two_images(prs, title, path1, path2, cap1="", cap2=""):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    left, top = Inches(0.5), Inches(0.3)
    tx = slide.shapes.add_textbox(left, top, Inches(9), Inches(0.5))
    tx.text_frame.text = title
    tx.text_frame.paragraphs[0].font.size = Pt(22)
    tx.text_frame.paragraphs[0].font.bold = True
    y = Inches(0.9)
    if os.path.isfile(path1):
        slide.shapes.add_picture(path1, Inches(0.5), y, width=Inches(4.4), height=Inches(3.2))
    if cap1:
        slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(4.4), Inches(0.35)).text_frame.text = cap1
    if os.path.isfile(path2):
        slide.shapes.add_picture(path2, Inches(5.0), y, width=Inches(4.4), height=Inches(3.2))
    if cap2:
        slide.shapes.add_textbox(Inches(5.0), Inches(4.3), Inches(4.4), Inches(0.35)).text_frame.text = cap2
    return slide


def add_slides_per_model_visuals(prs, report_dir, prefix, section_title, caption_note=""):
    """Add slides for all 7 models: 2 per slide, last slide 1 image. prefix e.g. 'confusion_matrix' or 'training_history'."""
    for i in range(0, len(MODELS), 2):
        pair = MODELS[i : i + 2]
        paths = [os.path.join(report_dir, f"{prefix}_{suffix}.png") for suffix, _ in pair]
        names = [name for _, name in pair]
        if len(pair) == 2:
            add_slide_two_images(
                prs,
                f"{section_title} — {names[0]} & {names[1]}",
                paths[0], paths[1],
                names[0], names[1],
            )
        else:
            if os.path.isfile(paths[0]):
                add_slide_with_image(
                    prs,
                    f"{section_title} — {names[0]}",
                    paths[0],
                    caption_note or names[0],
                )


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. Title
    add_title_slide(
        prs,
        "Lung Cancer Risk Prediction with Deep Learning",
        "Comparison with Dritsas & Trigka (2022) — Same dataset, neural networks vs classical ML"
    )

    # 2. Objective
    add_content_slide(prs, "Objective", [
        "Same task as the 2022 paper: predict lung cancer risk (YES/NO) from survey data.",
        "Same dataset: lung cancer survey (habits, symptoms, ~309 samples, 15 features).",
        "Their approach: classical ML (Rotation Forest) — 97.1% accuracy, 99.3% AUC.",
        "Our approach: deep learning (7 MLP architectures) — compare performance.",
    ])

    # 3. Data
    add_content_slide(prs, "Data", [
        "Source: survey lung cancer.csv (data/raw/).",
        "Features: age, gender, smoking, yellow fingers, anxiety, peer pressure, chronic disease,",
        "fatigue, allergy, wheezing, alcohol, coughing, shortness of breath, swallowing, chest pain.",
        "Target: LUNG_CANCER (YES/NO).",
        "Preprocessing: standard scaling, stratified splits, class weights for imbalance.",
    ])

    # 3b. Pipeline / Steps (high level)
    add_content_slide(prs, "Pipeline — steps we followed", [
        "1. Load data (survey lung cancer.csv) and encode labels (YES/NO → 1/0).",
        "2. Preprocess: StandardScaler (zero mean, unit variance), stratified train/test split.",
        "3. Train 7 architectures with the same hyperparameters (lr, batch, epochs, early stopping).",
        "4. Evaluate: single 80/20 split, 5-fold CV, 10-fold CV; report accuracy, F1, ROC-AUC.",
        "5. Compare with Dritsas & Trigka (2022): same dataset, different models (DL vs Rotation Forest).",
    ])

    # 4. Methods — 7 architectures
    add_content_slide(prs, "Methods — 7 deep learning architectures", [
        "Standard (128-64-32): 3-layer MLP, ReLU, BatchNorm, Dropout 0.3.",
        "Simple (64-32): 2-layer (64→32), ReLU, BatchNorm, Dropout 0.2.",
        "Minimal (32): 1 hidden layer (32 units), ReLU, BatchNorm, Dropout 0.2.",
        "Deep (5 layers): 5 hidden layers (64 each), ReLU, BatchNorm, Dropout 0.3.",
        "ELU Activation: same as Standard but ELU instead of ReLU.",
        "LayerNorm: same as Standard but LayerNorm instead of BatchNorm.",
        "Residual: skip connections, hidden size 128, Dropout 0.3.",
    ])

    # 5. Training setup (hyperparameters — same for all models)
    add_content_slide(prs, "Training setup (same for all architectures)", [
        "Learning rate: 0.001 (Adam optimizer).",
        "Batch size: 32.",
        "Max epochs: 100; early stopping if no improvement for 15 epochs.",
        "Loss: weighted Cross-Entropy (class weights for imbalanced YES/NO).",
        "Learning rate scheduler: ReduceLROnPlateau (factor 0.5, patience 10 epochs).",
        "Random seed: 42 (reproducibility).",
        "Input: 15 features, standardized (zero mean, unit variance).",
    ])

    # 6. Evaluation — what each setting means
    add_content_slide(prs, "Evaluation — what we report", [
        "Single 80/20 split: We split the data once — 80% for training, 20% held out for testing.",
        "  Train each model on the 80%, report accuracy on the 20%. One number per model.",
        "5-fold CV: Data split into 5 parts. 5 rounds: each round, 4 parts train, 1 part test.",
        "  Report mean and standard deviation of the 5 test accuracies. More stable.",
        "10-fold CV: Same idea with 10 parts. More rounds, often a bit more stable estimate.",
        "Metrics: accuracy, precision, recall, F1-score, ROC-AUC.",
    ])

    # 7. 5-fold vs 10-fold CV (short)
    add_content_slide(prs, "5-fold vs 10-fold CV — what is the difference?", [
        "5-fold: 5 train/test rounds; 10-fold: 10 rounds. Same idea: average over rounds.",
        "More folds → more stable mean and std; results are usually similar.",
        "We run both so you can compare and cite standard practice (both are common in papers).",
    ])

    # 8. Results — Single split (key figure)
    acc_single = os.path.join(SINGLE, "accuracy_comparison.png")
    add_slide_with_image(
        prs,
        "Results — Single 80/20 split",
        acc_single,
        "Test accuracy by model (red line = 2022 reference 97.1%)."
    )

    # 9. Results — Single split (complexity)
    comp_single = os.path.join(SINGLE, "complexity_vs_performance.png")
    add_slide_with_image(
        prs,
        "Results — Single split: complexity vs performance",
        comp_single,
        "Parameters vs test accuracy."
    )

    # 9b. Results — Metrics: Accuracy, Precision, Recall, F1 (single split)
    metrics_single = os.path.join(SINGLE, "comprehensive_metrics.png")
    add_slide_with_image(
        prs,
        "Results — Metrics: Accuracy, Precision, Recall, F1 (single split)",
        metrics_single,
        "All four metrics per model (single 80/20 run)."
    )

    # 9c. Confusion matrices — all 7 architectures (single split)
    add_slides_per_model_visuals(
        prs, SINGLE, "confusion_matrix",
        "Results — Confusion matrix",
        "True vs predicted (NO/YES). Diagonal = correct.",
    )

    # 9d. Training curves — all 7 architectures (single split)
    add_slides_per_model_visuals(
        prs, SINGLE, "training_history",
        "Results — Training curves",
        "Loss and accuracy over epochs; early stopping 15.",
    )

    # 9e. ROC curves — all 7 architectures (single split)
    add_slides_per_model_visuals(
        prs, SINGLE, "roc_curve",
        "Results — ROC curve",
        "True Positive Rate vs False Positive Rate; AUC in legend.",
    )

    # 9f. PR curves — all 7 architectures (single split)
    add_slides_per_model_visuals(
        prs, SINGLE, "pr_curve",
        "Results — Precision-Recall curve",
        "Precision vs Recall; AP in legend.",
    )

    # 10. Results — 5-fold CV
    acc_cv5 = os.path.join(CV5, "accuracy_comparison.png")
    add_slide_with_image(
        prs,
        "Results — 5-fold cross-validation (mean ± std)",
        acc_cv5,
        "More stable estimate; comparable to how the 2022 paper likely reported results."
    )

    # 11. Results — 5-fold vs 10-fold (if both exist)
    acc_cv10 = os.path.join(CV10, "accuracy_comparison.png")
    if os.path.isfile(acc_cv5) and os.path.isfile(acc_cv10):
        add_slide_two_images(
            prs,
            "Results — 5-fold vs 10-fold CV",
            acc_cv5, acc_cv10,
            "5-fold CV", "10-fold CV"
        )

    # 12. Comparison with 2022 paper
    add_content_slide(prs, "Comparison with Dritsas & Trigka (2022)", [
        "Reference: Rotation Forest — 97.1% accuracy, 99.3% AUC.",
        "Our best (e.g. LayerNorm with 5-fold CV): ~92% accuracy, ~94% AUC.",
        "Gap is expected: on small tabular data, tree ensembles often outperform neural nets.",
        "We use the same dataset and same task; we compare a different family of models (DL).",
        "Conclusion: DL is viable; classical ML remains strong on this dataset.",
    ])

    # 13. Summary / Takeaway
    add_content_slide(prs, "Summary", [
        "Same lung cancer survey data and classification task as the 2022 paper.",
        "Seven DL architectures trained and compared (single split + 5-fold and 10-fold CV).",
        "Training: lr 0.001, batch 32, early stopping 15; same setup for all models.",
        "Best DL performance ~92% accuracy (vs 97.1% Rotation Forest).",
        "All visuals and CSVs in reports/single_split, reports/cv5, reports/cv10.",
    ])

    # 14. Thank you / Questions
    add_title_slide(prs, "Thank you", "Questions?")

    os.makedirs(REPORTS, exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"Presentation saved to: {OUTPUT_PATH}")
    return OUTPUT_PATH


if __name__ == "__main__":
    main()
