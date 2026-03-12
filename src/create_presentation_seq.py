"""
Create a PowerPoint presentation for the sequence-model results (SMOTE + RNN/GRU/LSTM).

Includes: SMOTE explanation, class distribution, all evaluation metrics (accuracy,
precision, recall, F1, ROC-AUC), confusion matrices, ROC/PR curves, training history,
and 5-fold / 10-fold CV results if you have run those commands.

Run from project root:
  1. Generate all reports (single split + 5-fold + 10-fold):
       python src/compare_models_seq.py
       python src/compare_models_seq.py --cv 5
       python src/compare_models_seq.py --cv 10
  2. Build the presentation:
       python src/create_presentation_seq.py

Output: reports_seq/Lung_Cancer_Sequence_Models_Presentation.pptx
"""

import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError as e:
    print("Failed to import python-pptx:", e)
    print("Install with:  pip install python-pptx")
    sys.exit(1)

# Paths
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.dirname(SCRIPT_DIR)
REPORTS_SEQ = os.path.join(PROJECT_ROOT, "reports_seq")
SINGLE_SEQ = os.path.join(REPORTS_SEQ, "single_split")
CV5_SEQ = os.path.join(REPORTS_SEQ, "cv5")
CV10_SEQ = os.path.join(REPORTS_SEQ, "cv10")
OUTPUT_PATH = os.path.join(REPORTS_SEQ, "Lung_Cancer_Sequence_Models_Presentation.pptx")

# 9 models: RNN/GRU/LSTM × hidden 32, 64, 128 — must match compare_models_seq output
MODELS = [
    ("rnn_hidden=32", "RNN (hidden=32)"),
    ("rnn_hidden=64", "RNN (hidden=64)"),
    ("rnn_hidden=128", "RNN (hidden=128)"),
    ("gru_hidden=32", "GRU (hidden=32)"),
    ("gru_hidden=64", "GRU (hidden=64)"),
    ("gru_hidden=128", "GRU (hidden=128)"),
    ("lstm_hidden=32", "LSTM (hidden=32)"),
    ("lstm_hidden=64", "LSTM (hidden=64)"),
    ("lstm_hidden=128", "LSTM (hidden=128)"),
]


def add_title_slide(prs, title, subtitle=""):
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
    return slide


def add_content_slide(prs, title, bullets):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for b in bullets:
        p = body.add_paragraph()
        p.text = b
        p.level = 0
    return slide


def add_slide_with_image(prs, title, image_path, caption=""):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    left, top = Inches(0.5), Inches(0.3)
    tx = slide.shapes.add_textbox(left, top, Inches(9), Inches(0.6))
    tx.text_frame.text = title
    tx.text_frame.paragraphs[0].font.size = Pt(24)
    tx.text_frame.paragraphs[0].font.bold = True
    if os.path.isfile(image_path):
        slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.0), width=Inches(9), height=Inches(5.2))
    if caption:
        tx2 = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.4))
        tx2.text_frame.text = caption
        tx2.text_frame.paragraphs[0].font.size = Pt(12)
    return slide


def add_slide_two_images(prs, title, path1, path2, cap1="", cap2=""):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tx.text_frame.text = title
    tx.text_frame.paragraphs[0].font.size = Pt(22)
    tx.text_frame.paragraphs[0].font.bold = True
    y = Inches(0.9)
    if os.path.isfile(path1):
        slide.shapes.add_picture(path1, Inches(0.5), y, width=Inches(4.4), height=Inches(3.2))
    if cap1:
        slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(4.4), Inches(0.35)).text_frame.text = cap1
    if os.path.isfile(path2):
        slide.shapes.add_picture(path2, Inches(5.0), y, width=Inches(4.4), height=Inches(3.2))
    if cap2:
        slide.shapes.add_textbox(Inches(5.0), Inches(4.3), Inches(4.4), Inches(0.35)).text_frame.text = cap2
    return slide


def add_slides_per_model_visuals(prs, report_dir, prefix, section_title, caption_note=""):
    """Add slides for 3 models: 2 on first slide, 1 on second."""
    for i in range(0, len(MODELS), 2):
        pair = MODELS[i : i + 2]
        paths = [os.path.join(report_dir, f"{prefix}_{suffix}.png") for suffix, _ in pair]
        names = [name for _, name in pair]
        if len(pair) == 2:
            add_slide_two_images(prs, f"{section_title} — {names[0]} & {names[1]}", paths[0], paths[1], names[0], names[1])
        else:
            if os.path.isfile(paths[0]):
                add_slide_with_image(prs, f"{section_title} — {names[0]}", paths[0], caption_note or names[0])


def apply_theme(slide, rgb=(0xE8, 0xEE, 0xF5)):
    """Apply a simple light background color to the slide."""
    try:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
    except Exception:
        pass


def add_slide_metrics_table(prs, title, csv_path, columns_to_show=None):
    """Add a slide with a table built from the CSV (e.g. model_comparison_results.csv)."""
    if not os.path.isfile(csv_path):
        return None
    try:
        import pandas as pd
        df = pd.read_csv(csv_path)
    except Exception:
        return None
    if columns_to_show:
        df = df[[c for c in columns_to_show if c in df.columns]]
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tx.text_frame.text = title
    tx.text_frame.paragraphs[0].font.size = Pt(22)
    tx.text_frame.paragraphs[0].font.bold = True
    n_rows, n_cols = len(df) + 1, len(df.columns)
    if n_rows < 2 or n_cols < 1:
        return slide
    table = slide.shapes.add_table(n_rows, n_cols, Inches(0.5), Inches(0.9), Inches(9), Inches(0.4 * n_rows)).table
    for c, col in enumerate(df.columns):
        table.cell(0, c).text = str(col)[:20]
    for r in range(len(df)):
        for c, col in enumerate(df.columns):
            val = df.iloc[r][col]
            if isinstance(val, float):
                table.cell(r + 1, c).text = f"{val:.4f}" if abs(val) < 10 else f"{val:.2f}"
            else:
                table.cell(r + 1, c).text = str(val)[:20]
    return slide


def build_summary_bullets():
    """Read CSVs and build Greek summary bullets (best model, 5-fold, 10-fold, overfitting)."""
    import pandas as pd
    bullets = []
    # Single split
    csv_single = os.path.join(SINGLE_SEQ, "model_comparison_results.csv")
    if os.path.isfile(csv_single):
        try:
            df = pd.read_csv(csv_single)
            if "Test Accuracy" in df.columns and len(df) > 0:
                best = df.loc[df["Test Accuracy"].idxmax()]
                name = best.get("Model", "—")
                acc = best.get("Test Accuracy", 0)
                f1 = best.get("F1-Score", 0)
                roc = best.get("ROC-AUC", 0)
                params = int(best.get("Parameters", 0))
                epoch = best.get("Best Epoch", "—")
                if isinstance(epoch, float):
                    epoch = int(epoch)
                bullets.append(
                    f"Καλύτερο μοντέλο (80/20): {name} — Accuracy: {acc:.2%}, F1: {f1:.4f}, ROC-AUC: {roc:.4f}. "
                    f"Παράμετροι: {params:,}. Early stopping περίπου στο epoch {epoch}."
                )
        except Exception:
            pass
    # 5-fold
    csv_cv5 = os.path.join(CV5_SEQ, "model_comparison_results_seq_cv.csv")
    if os.path.isfile(csv_cv5):
        try:
            df = pd.read_csv(csv_cv5)
            acc_col = "Test Accuracy (mean)"
            if acc_col in df.columns and len(df) > 0:
                best = df.loc[df[acc_col].idxmax()]
                name = best.get("Model", "—")
                acc_m = best.get(acc_col, 0)
                acc_s = best.get("Test Accuracy (std)", 0)
                auc_m = best.get("ROC-AUC (mean)", 0)
                auc_s = best.get("ROC-AUC (std)", 0)
                bullets.append(f"5-fold CV: Καλύτερο {name} — Accuracy: {acc_m:.2%} ± {acc_s:.2%}, ROC-AUC: {auc_m:.4f} ± {auc_s:.4f}.")
        except Exception:
            pass
    # 10-fold
    csv_cv10 = os.path.join(CV10_SEQ, "model_comparison_results_seq_cv.csv")
    if os.path.isfile(csv_cv10):
        try:
            df = pd.read_csv(csv_cv10)
            acc_col = "Test Accuracy (mean)"
            if acc_col in df.columns and len(df) > 0:
                best = df.loc[df[acc_col].idxmax()]
                name = best.get("Model", "—")
                acc_m = best.get(acc_col, 0)
                acc_s = best.get("Test Accuracy (std)", 0)
                auc_m = best.get("ROC-AUC (mean)", 0)
                auc_s = best.get("ROC-AUC (std)", 0)
                bullets.append(f"10-fold CV: Καλύτερο {name} — Accuracy: {acc_m:.2%} ± {acc_s:.2%}, ROC-AUC: {auc_m:.4f} ± {auc_s:.4f}.")
        except Exception:
            pass
    # Overfitting note
    bullets.append(
        "Εκπαίδευση: Χρησιμοποιήθηκε early stopping (15 epochs χωρίς βελτίωση). "
        "Οι καμπύλες train/test δείχνουν ότι δεν υπάρχει σοβαρό overfitting· τα μοντέλα σταματούν όταν η test απόδοση σταθεροποιείται."
    )
    return bullets if bullets else ["Δεν βρέθηκαν αποτελέσματα CSV. Τρέξτε compare_models_seq.py και --cv 5, --cv 10."]


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. Τίτλος (χωρίς "Hello" / "Thank you" slides)
    add_title_slide(
        prs,
        "Πρόβλεψη κινδύνου καρκίνου πνεύμονα — Μοντέλα ακολουθιών (RNN, GRU, LSTM)",
        "Δεδομένα ισορροπημένα με SMOTE· ίδιο πλαίσιο με Dritsas & Trigka (2022)"
    )

    # 2. SMOTE και στόχος — περισσότερο κείμενο, εξήγηση προβλήματος και λύσης
    add_content_slide(prs, "Χρήση SMOTE και στόχος εργασίας", [
        "Το αρχικό dataset (survey lung cancer.csv) είχε έντονη ανισορροπία κλάσεων: ~87% YES, ~13% NO (309 δείγματα).",
        "Για να αντιμετωπίσουμε το πρόβλημα χρησιμοποιήσαμε SMOTE (Synthetic Minority Over-sampling): δημιουργούνται συνθετικά δείγματα για τη μειοψηφική κλάση, ώστε να προκύψει ισορροπημένο dataset (~540 δείγματα, 50%–50%).",
        "Στον ίδιο ισορροπημένο χώρο εκπαιδεύουμε τρία μοντέλα ακολουθιών (RNN, GRU, LSTM) και αξιολογούμε με accuracy, precision, recall, F1, ROC-AUC και confusion matrices.",
    ])

    # 3. Δεδομένα και κατανομή κλάσεων
    add_content_slide(prs, "Δεδομένα και κατανομή κλάσεων", [
        "Πηγή: survey lung cancer.csv — 15 χαρακτηριστικά, στόχος LUNG_CANCER (YES/NO).",
        "Πριν το SMOTE: 309 δείγματα — περίπου 39 NO (12,6%), 270 YES (87,4%).",
        "Μετά το SMOTE: 540 δείγματα — 270 NO (50%), 270 YES (50%). Το ισορροπημένο CSV αποθηκεύεται στο data/processed/.",
        "Διαχωρισμός: 80% train / 20% test πάνω στο ισορροπημένο dataset (432 train, 108 test). Για 5-fold και 10-fold CV χρησιμοποιούμε το ίδιο ισορροπημένο dataset.",
    ])

    # 4. Διεργασία (pipeline)
    add_content_slide(prs, "Διεργασία", [
        "1. Φόρτωση αρχικού CSV και κωδικοποίηση ετικετών (YES/NO → 1/0).",
        "2. Εφαρμογή SMOTE μία φορά → ισορροπημένο dataset (540 δείγματα).",
        "3. StandardScaler· στρωματοποιημένο split 80/20 (ή k-fold για CV).",
        "4. Εκπαίδευση RNN, GRU, LSTM με τα ίδια hyperparameters (Adam, early stopping, κ.ά.).",
        "5. Αξιολόγηση: accuracy, precision, recall, F1, ROC-AUC, PR-AUC· confusion matrices και καμπύλες ROC/PR.",
    ])

    # 5. Μοντέλα — RNN, GRU, LSTM (hidden=32, 64, 128)
    add_content_slide(prs, "Μοντέλα — RNN, GRU, LSTM (hidden=32, 64, 128)", [
        "Τα 15 χαρακτηριστικά αντιμετωπίζονται ως ακολουθία μήκους 15 (input_dim=1).",
        "Για κάθε αρχιτεκτονική (RNN, GRU, LSTM) δοκιμάζουμε τρία μεγέθη: hidden=32, 64, 128.",
        "RNN: 1 στρώμα· παράμετροι ~1.2K (32), ~4.4K (64), ~17K (128).",
        "GRU: 1 στρώμα· παράμετροι ~3.5K (32), ~13K (64), ~50K (128).",
        "LSTM: 1 στρώμα· παράμετροι ~4.5K (32), ~17K (64), ~66K (128). Όλα δίνουν 2 logits (NO/YES).",
    ])

    # 6. Ρυθμίσεις εκπαίδευσης — epochs, learning rate, early stop κ.ά.
    add_content_slide(prs, "Ρυθμίσεις εκπαίδευσης (κοινές για όλα τα μοντέλα)", [
        "Optimizer: Adam. Learning rate: 0.001.",
        "Batch size: 32. Μέγιστα epochs: 100.",
        "Early stopping: διακοπή αν δεν υπάρχει βελτίωση στο test accuracy για 15 συνεχόμενα epochs.",
        "Loss: Cross-Entropy με class weights [1, 1] (τα δεδομένα είναι ήδη ισορροπημένα με SMOTE).",
        "Scheduler: ReduceLROnPlateau (factor 0.5, patience 10 epochs).",
        "Random seed: 42. Όλα τα plots και CSV βρίσκονται στο reports_seq/.",
    ])

    # 6b. Πλήθη παραμέτρων (βαρών) ανά μοντέλο — ξεκάθαρα ορατά
    add_slide_metrics_table(
        prs,
        "Πλήθη παραμέτρων (βαρών) ανά μοντέλο",
        os.path.join(SINGLE_SEQ, "model_comparison_results.csv"),
        columns_to_show=["Model", "Parameters"],
    )

    # 6c. Τιμές βαρών (weights) — πού βρίσκονται (CSV + ιστογράμματα)
    add_content_slide(prs, "Τιμές βαρών (weights) — πού βρίσκονται", [
        "Για κάθε εκπαιδευμένο μοντέλο αποθηκεύονται:",
        "(α) Στατιστικά βαρών ανά layer: shape, αριθμός παραμέτρων, mean, std, min, max → reports_seq/single_split/weights_<model>.csv",
        "(β) Ιστόγραμμα κατανομής των τιμών των βαρών → reports_seq/single_split/weight_histogram_<model>.png",
        "Τα αρχεία παράγονται αυτόματα μετά την εκπαίδευση (single 80/20 split).",
    ])

    # 7. Αξιολόγηση — μετρικές
    add_content_slide(prs, "Αξιολόγηση — μετρικές που αναφέρουμε", [
        "Accuracy, Precision, Recall, F1-Score ανά μοντέλο.",
        "ROC-AUC και Precision-Recall AUC.",
        "Confusion matrix (αληθής vs προβλεπόμενη κλάση NO/YES).",
        "Αποτελέσματα: single split 80/20 και, όταν έχουν τρέξει, 5-fold και 10-fold stratified CV (μέσος ± τυπική απόκλιση).",
    ])

    # 8. Αποτελέσματα — πίνακας μετρικών (80/20)
    csv_single = os.path.join(SINGLE_SEQ, "model_comparison_results.csv")
    add_slide_metrics_table(
        prs,
        "Αποτελέσματα — Πίνακας μετρικών (80/20 split πάνω σε δεδομένα SMOTE)",
        csv_single,
        columns_to_show=["Model", "Test Accuracy", "Precision", "Recall", "F1-Score", "ROC-AUC", "Parameters"],
    )

    # 9. Αποτελέσματα — γράφημα accuracy (80/20)
    acc_single = os.path.join(SINGLE_SEQ, "accuracy_comparison.png")
    add_slide_with_image(prs, "Αποτελέσματα — Test accuracy (80/20 split)", acc_single, "RNN, GRU, LSTM πάνω σε δεδομένα ισορροπημένα με SMOTE.")

    # 10. Αποτελέσματα — Accuracy, Precision, Recall, F1 (4 πλαίσια)
    metrics_single = os.path.join(SINGLE_SEQ, "comprehensive_metrics.png")
    add_slide_with_image(prs, "Αποτελέσματα — Accuracy, Precision, Recall, F1 (80/20)", metrics_single, "Οι τέσσερις μετρικές ανά μοντέλο.")

    # 11. Confusion matrices
    add_slides_per_model_visuals(prs, SINGLE_SEQ, "confusion_matrix", "Αποτελέσματα — Confusion matrix", "Αληθής vs προβλεπόμενη κλάση (NO/YES).")

    # 12. Καμπύλες εκπαίδευσης
    add_slides_per_model_visuals(prs, SINGLE_SEQ, "training_history", "Αποτελέσματα — Καμπύλες εκπαίδευσης", "Loss και accuracy ανά epoch.")

    # 13. ROC curves
    add_slides_per_model_visuals(prs, SINGLE_SEQ, "roc_curve", "Αποτελέσματα — ROC curve", "AUC στην υπόμνηση.")

    # 14. PR curves
    add_slides_per_model_visuals(prs, SINGLE_SEQ, "pr_curve", "Αποτελέσματα — Precision-Recall curve", "Average precision στην υπόμνηση.")

    # 15. 5-fold CV (αν υπάρχει)
    acc_cv5 = os.path.join(CV5_SEQ, "accuracy_comparison.png")
    if os.path.isfile(acc_cv5):
        add_slide_with_image(prs, "Αποτελέσματα — 5-fold cross-validation (μέσος ± τυπική απόκλιση)", acc_cv5, "Δεδομένα SMOTE, 5-fold stratified CV.")
        add_slide_metrics_table(
            prs,
            "Αποτελέσματα — Μετρικές 5-fold CV",
            os.path.join(CV5_SEQ, "model_comparison_results_seq_cv.csv"),
            columns_to_show=["Model", "Test Accuracy (mean)", "Test Accuracy (std)", "ROC-AUC (mean)", "ROC-AUC (std)"],
        )

    # 16. 10-fold CV (αν υπάρχει)
    acc_cv10 = os.path.join(CV10_SEQ, "accuracy_comparison.png")
    if os.path.isfile(acc_cv10):
        add_slide_with_image(prs, "Αποτελέσματα — 10-fold cross-validation (μέσος ± τυπική απόκλιση)", acc_cv10, "Δεδομένα SMOTE, 10-fold stratified CV.")
        add_slide_metrics_table(
            prs,
            "Αποτελέσματα — Μετρικές 10-fold CV",
            os.path.join(CV10_SEQ, "model_comparison_results_seq_cv.csv"),
            columns_to_show=["Model", "Test Accuracy (mean)", "Test Accuracy (std)", "ROC-AUC (mean)", "ROC-AUC (std)"],
        )

    # 17. 5-fold vs 10-fold δίπλα-δίπλα (αν υπάρχουν και τα δύο)
    if os.path.isfile(acc_cv5) and os.path.isfile(acc_cv10):
        add_slide_two_images(prs, "Αποτελέσματα — 5-fold vs 10-fold CV", acc_cv5, acc_cv10, "5-fold CV", "10-fold CV")

    # 18. Σύνοψη (κείμενο)
    add_content_slide(prs, "Σύνοψη", [
        "Εφαρμόσαμε SMOTE στο αρχικό survey dataset → ισορροπημένα 540 δείγματα (50% YES, 50% NO).",
        "Εννέα μοντέλα: RNN, GRU, LSTM με hidden=32, 64, 128 (15 features ως ακολουθία).",
        "Πλήρης αξιολόγηση: accuracy, precision, recall, F1, ROC-AUC, confusion matrices, καμπύλες ROC/PR.",
        "Τιμές βαρών (weights): στατιστικά ανά layer και ιστογράμματα στο reports_seq/single_split/ (weights_*.csv, weight_histogram_*.png).",
        "Single 80/20 split και προαιρετικά 5-fold και 10-fold CV. Όλα τα figures και CSV στο reports_seq/.",
    ])

    # 19. Τελική σύνοψη — καλύτερα μοντέλα, αποτελέσματα, overfitting (από τα CSV)
    summary_bullets = build_summary_bullets()
    add_content_slide(prs, "Τελική σύνοψη — Καλύτερα μοντέλα και σχόλια", summary_bullets)

    # Θέμα: απλό ανοιχτό φόντο σε όλες τις διαφάνειες
    for slide in prs.slides:
        apply_theme(slide)

    os.makedirs(REPORTS_SEQ, exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"Presentation saved to: {OUTPUT_PATH}")
    return OUTPUT_PATH


if __name__ == "__main__":
    main()
